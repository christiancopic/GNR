# Author: Christian Copic, 2024

import os
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import math
import csv
import time
import numpy as np


def main():
    start = time.time()
    cwd = os.getcwd()
    cwd = cwd
    completed = []
    # Adding count for debugging purposes
    count = 0
    total = 0
    # Will have to change this line to be compatible with windows
    for file in os.listdir(cwd + '/devicedataset/Before'):
        # DEBUG
        if count < 1000000:
            if file[-4:] == '.csv':

                # File structure is a mess... will be using quick and dirty
                # method to distinguish them... hope it works well :)
                
                device_name = file[:3].lower()
                vds = file[3:5]
                device_count_neg = 0
                device_count_pos = 0
                file_names = []
                flag = False
                
                # DEBUG
                #if device_name == '58b':
                    #print('e')

                for file_after in os.listdir(cwd + '/devicedataset/After'):
                    # Ugly code :(
                    # DEBUG
                    #if '58b'  in file_after and device_name == '58b':
                        #print('e')

                    if device_name in file_after.lower() and vds in file_after and vds == '-1':
                        if(len(file_after) > 16):
                            continue
                        device_count_neg = device_count_neg + 1
                        total = total + device_count_neg
                        #flag = True
                        process_ald_csvs(file, file_after, cwd, file_names, device_name, True, True)
                    if device_name in file_after.lower() and vds in file_after and vds == '+1':
                        if(len(file_after) > 16):
                            continue
                        device_count_pos = device_count_pos + 1
                        total = total + device_count_pos
                        #flag = True
                        process_ald_csvs(file, file_after, cwd, file_names, device_name, True, True)

                if(device_count_neg > 1 or device_count_pos > 1):
                    print('Error: Too many files for device ' + device_name + '. Delete extra csv.')
                    exit(0)

                #print(device_name + ' has ' + str(device_count) + ' hits in After')

                if False: 
                    data_before = pd.read_csv(cwd + '/devicedataset/Before/' + file)
                    data_after = pd.read_csv(cwd + '/devicedataset/After/' + file_after)
                    # Fix name issue
                    data_before = data_before.rename(columns={'Voltage Amplitude Setpoint (V).1':'Vg'})
                    data_after = data_after.rename(columns={'Voltage Amplitude Setpoint (V).1':'Vg'})
                    file_names.append(file)
                    file_names.append(file_after)
                    make_graph(data_before, data_after,device_name, file_names)
                flag = False
                count = count + 1
        else:
            continue
    create_pptx()
    makeCDFs()
    print('Done. Total time: ' + str(time.time() - start))

def process_ald_csvs(file, file_after, cwd, file_names, device_name, makeGraph, analyze):
    data_before = pd.read_csv(cwd + '/devicedataset/Before/' + file)
    data_after = pd.read_csv(cwd + '/devicedataset/After/' + file_after)
    # Fix name issue
    data_before = data_before.rename(columns={'Voltage Amplitude Setpoint (V).1':'Vg'})
    data_after = data_after.rename(columns={'Voltage Amplitude Setpoint (V).1':'Vg'})
    file_names.append(file)
    file_names.append(file_after)
    if(device_name == '4af'):
        print('debug')
    if(len(file_names) > 2):
        print('debug')
    if makeGraph and False:
        make_graph(data_before, data_after,device_name, file_names)
    if analyze and False:
        e_char(data_before, data_after,device_name, file_names)

# Create ID vs. VGS curves of data
def make_graph(data_before, data_after, device_name, file_names):
    #log y with regular x axis
    plt.semilogy(figsize=(9, 5))

    data_before['Id'] = abs(data_before['Id'])
    data_after['Id'] = abs(data_after['Id'])
    # Assumption that same Vd values exist in data_after
    VD_values = data_before['Vd'].unique()
    if(VD_values[0] < 0):
        vd = '-1'
    else:
        vd = '+1'

    for i in range(len(VD_values)):
        subdf_before = data_before[data_before['Vd'] == VD_values[i]]
        subdf_after = data_after[data_after['Vd'] == VD_values[i]]
        plt.plot(subdf_before['Vg'], subdf_before['Id'], label = '$V_{DS}$ = ' + str(VD_values[i]) + 'V Before ALD')
        plt.plot(subdf_after['Vg'], subdf_after['Id'], label = '$V_{DS}$ = ' + str(VD_values[i]) + 'V After ALD')
    
    # Use LaTeX for labels
    plt.xlabel('$V_{GS}$ (V)')
    plt.ylabel('$-I_{D}$ (A)')
    plt.title('$-I_{D}$ (A) vs. $V_{GS}$ (V) from ' + device_name)
    plt.legend()
    plt.savefig(device_name + vd + '.png')
    plt.close()

    #save to into pptx file code here:

# Calculate electrical characteristics of device
def e_char(data_before, data_after, device_name, file_names):
    edata_before = [] # Figure out better method but this works in mean time
    edata_after = []
    VD_values = data_before['Vd'].unique()
    for i in range(len(VD_values)):
        subdf_before = data_before[data_before['Vd'] == VD_values[i]]
        subdf_after = data_after[data_after['Vd'] == VD_values[i]]
        # Reverse & Forward Sweep
        RS_before = subdf_before.iloc[60:,:]
        FS_before = subdf_before.iloc[:60,:]
        RS_after = subdf_after.iloc[60:,:]
        FS_after = subdf_after.iloc[:60,:]

        #print(device_name)
        #print(VD_values)
        #print(file_names)

        x = forward_reverse_analysis(RS_before, FS_before, VD_values, i)
        edata_before = edata_before + x
        y = forward_reverse_analysis(RS_after, FS_after, VD_values, i)
        edata_after = edata_after + y
        

    make_csv(edata_before, device_name, 'before', VD_values)
    make_csv(edata_after, device_name, 'after', VD_values)
    #put into csv/xlsx file code here:
    #also add maximum IG value for -1V = VDS


def forward_reverse_analysis(RS, FS, VD_values, i):
    edata = []

    for j in [RS,FS]:
            # Supress noise
            pd.options.mode.chained_assignment = None
            j['moving'] = j.rolling(5, center=True)['Id'].mean()
            #print(j['moving'])
            # high = about 0V
            high_df = j.iloc[(j['Vg']-0).abs().argsort()[:1]]
            # low = about -0.5V

            '''
            FIX THIS FOR BOTH VDS VALUES
            '''

            if(VD_values[0]<0):
                low_df = j.iloc[(j['Vg']+0.5).abs().argsort()[:1]]
            else:
                low_df = j.iloc[(j['Vg']-0.5).abs().argsort()[:1]]

            high_vg = float(high_df['Vg'].to_string(index=False, header=False))
            low_vg = float(low_df['Vg'].to_string(index=False, header=False))
            
            high_id = high_df.values[0][-1]
            low_id = low_df.values[0][-1]

            delta_vg = (high_vg-low_vg) * 1000 #mV
            delta_id = math.log10(high_id) - math.log10(low_id)
            SS = delta_vg/delta_id * -1
            #print('SUBTHRESHOLD for VDS= ' + str(VD_values[i]) + ' : ' + str(SS))

            edata.append(SS)
            edata.append(j['moving'].min(0))
            edata.append(j['moving'].max(0))
            edata.append(j['moving'].max(0)/j['moving'].min(0))
            #print(str(j))
            #print('For VD = ' + str(VD_values[i]) + ' on device ' + device_name + ': ')
            #print('Min ID Value: ' + str(j['ID'].min(0)))
            #print('Max ID Value: ' + str(j['ID'].max(0)))
            #print('Ion/Ioff ratio: ' + str(j['ID'].max(0) / subdf['ID'].min(0)))
            if(abs(VD_values[i]) == 1.0): 
                #print('Max IG Value: ' + str(j['IG'].max(0)))
                edata.append(j['Ig'].max(0))
    return edata



def create_pptx():
    ppt = Presentation()

    first_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = "Parsed Data Matplotlib Graphs"
    first_slide.shapes[0].text_frame.paragraphs[0].text = title
    
    for file in os.listdir(os.getcwd()):
        if file[-4:] == '.png':
            img = file
            slide = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(slide)
            slide.shapes.add_picture(img, left= Inches(2),top = Inches(1),height = Inches(5))
            os.remove(file)
 
    ppt.save('data_parser.pptx')

def make_csv(data, device_name, bef_aft, Vd):
    cwd = os.getcwd()
    if not os.path.exists(cwd + '/excelFiles'):
        os.makedirs(cwd + '/excelFiles')

    if (Vd[0] > 0):
        vd_value = '+1'
    else:
        vd_value = '-1'

    # This is ugly :(
    if vd_value == '+1':
        rows = [[str(device_name), 'Forward Sweep', '', 'Reverse Sweep', ''], 
            ['',           '0.1V',             '1V',                '0.1V',            '1V'], 
            ['Ion',        str(data[16]),       str(data[7]),        str(data[12]),      str(data[2])], 
            ['Ioff',       str(data[15]),       str(data[6]),        str(data[11]),      str(data[1])], 
            ['Ion/Ioff',   str(data[17]),       str(data[8]),        str(data[13]),      str(data[3])], 
            ['SS',         str(data[14]),       str(data[5]),        str(data[10]),      str(data[0])],
            ['Max Ig',     'n/a',              str(data[9]),        'n/a',             str(data[4])]]
        #print('DATA: ' + str(data))
    else:
        rows = [[str(device_name), 'Forward Sweep', '', 'Reverse Sweep', ''], 
            ['',           '-0.1V',             '-1V',                '-0.1V',            '-1V'], 
            ['Ion',        str(data[16]),       str(data[7]),        str(data[12]),      str(data[2])], 
            ['Ioff',       str(data[15]),       str(data[6]),        str(data[11]),      str(data[1])], 
            ['Ion/Ioff',   str(data[17]),       str(data[8]),        str(data[13]),      str(data[3])], 
            ['SS',         str(data[14]),       str(data[5]),        str(data[10]),      str(data[0])],
            ['Max Ig',     'n/a',              str(data[9]),        'n/a',             str(data[4])]]

    with open(cwd + '/excelFiles/' + device_name + vd_value + bef_aft + '.csv', 'w') as csvfile: 
        csvwriter = csv.writer(csvfile) 
        csvwriter.writerows(rows)

def makeCDFs():
    # [[Ion],[Ioff],[Ion/Ioff], [SS]]
    before_data = [[],[],[],[]]
    after_data = [[],[],[],[]]
    cwd = os.getcwd()

    for file in os.listdir(cwd + '/excelFiles/'):
        if file[-4:] == '.csv':
            if 'before' in file:
                df = pd.read_csv(cwd + '/excelFiles/' + file, skiprows=1)
                if '-1' in file:
                    vd = '-1'
                else:
                    vd = '1'
                
                if(df[vd + 'V'][2] < 50 or df[vd + 'V.1'][2] < 50):
                    continue

                before_data[0].append(df[vd + 'V'][0])
                before_data[1].append(df[vd + 'V'][1])
                before_data[2].append(df[vd + 'V'][2])
                before_data[3].append(df[vd + 'V'][3])

                before_data[0].append(df[vd + 'V.1'][0])
                before_data[1].append(df[vd + 'V.1'][1])
                before_data[2].append(df[vd + 'V.1'][2])
                before_data[3].append(df[vd + 'V.1'][3])

            if 'after' in file:
                df = pd.read_csv(cwd + '/excelFiles/' + file, skiprows=1)
                if '-1' in file:
                    vd = '-1'
                else:
                    vd = '1'
                
                if(df[vd + 'V'][2] < 50 or df[vd + 'V.1'][2] < 50):
                    continue

                after_data[0].append(df[vd + 'V'][0])
                after_data[1].append(df[vd + 'V'][1])
                after_data[2].append(df[vd + 'V'][2])
                after_data[3].append(df[vd + 'V'][3])

                after_data[0].append(df[vd + 'V.1'][0])
                after_data[1].append(df[vd + 'V.1'][1])
                after_data[2].append(df[vd + 'V.1'][2])
                after_data[3].append(df[vd + 'V.1'][3])

    IoffCDF(before_data[0], after_data[0])
    IonCDF(before_data[1], after_data[1])
    IRatioCDF(before_data[2], after_data[2])
    ssCDF(before_data[3], after_data[3])

def IoffCDF(before, after):
    for i in [before, after]:
        N = len(i)
        x = np.sort(i) 
        y = np.arange(N) / float(N) 
        plt.semilogx(x, y, marker='o') 

        #plt.hist(abs(df['Mobility']), density=True, cumulative=True, label='CDF',
        #    histtype='step', alpha=0.8, color='k')
    plt.legend(['before ALD', 'after ALD'])
    plt.xlabel('$I_{off}$')
    plt.ylabel('CDF')
    plt.title('$I_{off}$ CDF')
    plt.show()

def IonCDF(before, after):
    for i in [before, after]:
        N = len(i)
        x = np.sort(i) 
        y = np.arange(N) / float(N) 
        plt.semilogx(x, y, marker='o') 

        #plt.hist(abs(df['Mobility']), density=True, cumulative=True, label='CDF',
        #    histtype='step', alpha=0.8, color='k')
    plt.legend(['before ALD', 'after ALD'])
    plt.xlabel('$I_{on}$')
    plt.ylabel('CDF')
    plt.title('$I_{on}$ CDF')
    plt.show()

def IRatioCDF(before, after):
    for i in [before, after]:
        N = len(i)
        x = np.sort(i) 
        y = np.arange(N) / float(N) 
        plt.semilogx(x, y, marker='o') 

        #plt.hist(abs(df['Mobility']), density=True, cumulative=True, label='CDF',
        #    histtype='step', alpha=0.8, color='k')
    plt.legend(['before ALD', 'after ALD'])
    plt.xlabel('$I_{on}$/$I_{off}$')
    plt.ylabel('CDF')
    plt.title('$I_{on}$/$I_{off}$ CDF')
    plt.show()

def ssCDF(before, after):
    for i in [before, after]:
        N = len(i)
        x = np.sort(i) 
        y = np.arange(N) / float(N) 
        plt.semilogx(x, y, marker='o') 

        #plt.hist(abs(df['Mobility']), density=True, cumulative=True, label='CDF',
        #    histtype='step', alpha=0.8, color='k')
    plt.legend(['before ALD', 'after ALD'])
    plt.xlabel('Subthreshold Swing')
    plt.ylabel('CDF')
    plt.title('Subthreshold Swing CDF')
    plt.show()

if __name__ == '__main__':
    main()