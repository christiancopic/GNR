# Purpose: 
# Author: Christian Copic, 2023

import os
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import math
import csv
import time


def main():
    start = time.time()
    cwd = os.getcwd()
    # Adding count for debugging purposes
    count = 0
    # Will have to change this line to be compatible with windows
    for file in os.listdir(cwd + '/devicedataset'):
        if count < 1:
            if file[-4:] == '.csv':
                data = pd.read_csv(cwd + '/devicedataset/' + file)
                device_name = file.split('[')[1].split(']')[0].split(' ')[0]
                make_graph(data, device_name)
                e_char(data, device_name)

                count = count + 1
        else:
            continue
    create_pptx()
    print('Done. Total time: ' + str(time.time() - start))

# Create ID vs. VGS curves of data
def make_graph(data, device_name):
    #log y with regular x axis
    plt.semilogy(figsize=(9, 5))

    data['ID'] = abs(data['ID'])
    VD_values = data['VD'].unique()

    for i in range(len(VD_values)):
        subdf = data[data['VD'] == VD_values[i]]
        plt.plot(subdf['VG'], subdf['ID'], label = '$V_{DS}$ = ' + str(VD_values[i]) + 'V')
    
    # Use LaTeX for labels
    plt.xlabel('$V_{GS}$ (V)')
    plt.ylabel('$-I_{D}$ (A)')
    plt.title('$-I_{D}$ (A) vs. $V_{GS}$ (V) from ' + device_name)
    plt.legend()
    plt.savefig(device_name + '.png')
    plt.close()

    #save to into pptx file code here:

# Calculate electrical characteristics of device
def e_char(data, device_name):
    edata=[] # Figure out better method but this works in mean time
    VD_values = data['VD'].unique()
    for i in range(len(VD_values)):
        subdf = data[data['VD'] == VD_values[i]]
        # Reverse & Forward Sweep
        RS = subdf.iloc[300:,:]
        FS = subdf.iloc[:300,:]
        for j in [RS,FS]:
            # Supress noise
            pd.options.mode.chained_assignment = None
            j['moving'] = j.rolling(5, center=True)['ID'].mean()
            # high = about 0V
            high_df = j.iloc[(j['VG']-0).abs().argsort()[:1]]
            # low = about -0.5V
            low_df = j.iloc[(j['VG']+0.5).abs().argsort()[:1]]

            high_vg = float(high_df['VG'].to_string(index=False, header=False))
            low_vg = float(low_df['VG'].to_string(index=False, header=False))
            
            high_id = high_df.values[0][-1]
            low_id = low_df.values[0][-1]

            delta_vg = (high_vg-low_vg) * 1000 #mV
            delta_id = math.log10(high_id) - math.log10(low_id)
            SS = delta_vg/delta_id * -1
            #print('SUBTHRESHOLD for VDS= ' + str(VD_values[i]) + ' : ' + str(SS))

            edata.append(SS)
            edata.append(j['moving'].min(0))
            edata.append(j['moving'].max(0))
            edata.append(j['moving'].max(0)/j['moving'].min(0))
            #print(str(j))
            #print('For VD = ' + str(VD_values[i]) + ' on device ' + device_name + ': ')
            #print('Min ID Value: ' + str(j['ID'].min(0)))
            #print('Max ID Value: ' + str(j['ID'].max(0)))
            #print('Ion/Ioff ratio: ' + str(j['ID'].max(0) / subdf['ID'].min(0)))
            if(VD_values[i] == -1.0): 
                #print('Max IG Value: ' + str(j['IG'].max(0)))
                edata.append(j['IG'].max(0))

    make_csv(edata, device_name)
    #put into csv/xlsx file code here:
    #also add maximum IG value for -1V = VDS

def create_pptx():
    ppt = Presentation()

    first_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = "Parsed Data Matplotlib Graphs"
    first_slide.shapes[0].text_frame.paragraphs[0].text = title
    
    for file in os.listdir(os.getcwd()):
        if file[-4:] == '.png':
            img = file
            slide = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(slide)
            slide.shapes.add_picture(img, left= Inches(2),top = Inches(1),height = Inches(5))
            os.remove(file)
 
    ppt.save('data_parser.pptx')

def make_csv(data, device_name):
    cwd = os.getcwd()
    if not os.path.exists(cwd + '/excelFiles'):
        os.makedirs(cwd + '/excelFiles')

    # This is ugly :(
    rows = [[str(device_name), 'Forward Sweep', '', 'Reverse Sweep', ''], 
         ['',           '0.1V',             '1V',                '0.1V',            '1V'], 
         ['Ion',        str(data[6]),       str(data[15]),        str(data[2]),      str(data[10])], 
         ['Ioff',       str(data[5]),       str(data[14]),        str(data[1]),      str(data[9])], 
         ['Ion/Ioff',   str(data[7]),       str(data[16]),        str(data[3]),      str(data[11])], 
         ['SS',         str(data[4]),       str(data[13]),        str(data[0]),      str(data[8])],
         ['Max Ig',     'n/a',              str(data[17]),        'n/a',             str(data[12])]]
    #print('DATA: ' + str(data))

    with open(cwd + '/excelFiles/' + device_name + '.csv', 'w') as csvfile: 
        csvwriter = csv.writer(csvfile) 
        csvwriter.writerows(rows)

if __name__ == '__main__':
    main()